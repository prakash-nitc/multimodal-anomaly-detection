# -*- coding: utf-8 -*-
"""Build the Phase 2 review deck for DA-ZVAD.

Design intent: a panel review, not a conference talk. Slides are read while
being spoken over, so each carries one claim, states its evidence, and stops.
Numbers are large where they are the point. The narrative deliberately includes
the failed first run -- for a panel questioning whether the work is genuinely
the student's, a debugging story is stronger evidence than a clean table.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---------------------------------------------------------------- palette
INK        = RGBColor(0x14, 0x1C, 0x1E)
INK_2      = RGBColor(0x3A, 0x4A, 0x4C)
MUTED      = RGBColor(0x61, 0x72, 0x74)
ACCENT     = RGBColor(0x0D, 0x6B, 0x67)
ACCENT_LT  = RGBColor(0xDD, 0xED, 0xEB)
CAUTION    = RGBColor(0x8A, 0x57, 0x15)
CAUTION_LT = RGBColor(0xF6, 0xEB, 0xD8)
FAIL       = RGBColor(0x8C, 0x3A, 0x31)
FAIL_LT    = RGBColor(0xF6, 0xE4, 0xE1)
RULE       = RGBColor(0xD8, 0xE0, 0xDF)
SURF       = RGBColor(0xF5, 0xF7, 0xF7)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Segoe UI"
MONO = "Consolas"

W, H = 13.333, 7.5
L, R = 0.85, 0.85
CW = W - L - R          # content width

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]

_n = {"i": 0}


# ---------------------------------------------------------------- helpers
def _tf(shape, text, size, bold=False, color=INK, font=FONT, align=PP_ALIGN.LEFT,
        spacing=1.0, italic=False):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = spacing
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tf


def box(slide, x, y, w, h, text, size, **kw):
    s = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _tf(s, text, size, **kw)
    return s


def rect(slide, x, y, w, h, fill, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                               Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(0.75)
    s.shadow.inherit = False
    return s


def slide(title, eyebrow=None, number=True):
    """Standard content slide: eyebrow, title, accent rule."""
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, W, H, WHITE)
    y = 0.52
    if eyebrow:
        box(s, L, y, CW, 0.25, eyebrow.upper(), 11, bold=True, color=ACCENT)
        y += 0.30
    box(s, L, y, CW, 0.6, title, 27, bold=True, color=INK)
    rect(s, L, y + 0.72, 1.5, 0.035, ACCENT)
    if number:
        _n["i"] += 1
        box(s, W - R - 0.6, H - 0.52, 0.6, 0.25, str(_n["i"]), 10,
            color=MUTED, align=PP_ALIGN.RIGHT)
    return s


def bullets(slide_, items, top, left=L, width=CW, size=16, gap=0.46,
            color=INK_2, marker=True):
    y = top
    for it in items:
        if isinstance(it, tuple):
            head, rest = it
        else:
            head, rest = None, it
        if marker:
            rect(slide_, left, y + 0.10, 0.075, 0.075, ACCENT)
        s = slide_.shapes.add_textbox(Inches(left + (0.26 if marker else 0)),
                                      Inches(y - 0.04), Inches(width - 0.26),
                                      Inches(0.4))
        tf = s.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.line_spacing = 1.25
        if head:
            r = p.add_run(); r.text = head + "  "
            r.font.size = Pt(size); r.font.bold = True
            r.font.color.rgb = INK; r.font.name = FONT
        r = p.add_run(); r.text = rest
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = FONT
        y += gap
    return y


def table(slide_, data, x, y, w, col_w=None, size=13, header=True,
          hi_rows=(), row_h=0.42, num_cols=()):
    rows, cols = len(data), len(data[0])
    shp = slide_.shapes.add_table(rows, cols, Inches(x), Inches(y),
                                  Inches(w), Inches(row_h * rows))
    tbl = shp.table
    tbl.first_row = header
    tbl.horz_banding = False
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = Inches(cw)
    for ri, row in enumerate(data):
        tbl.rows[ri].height = Inches(row_h)
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = Inches(0.14)
            cell.margin_right = Inches(0.14)
            cell.margin_top = Inches(0.05)
            cell.margin_bottom = Inches(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if ri == 0 and header:
                cell.fill.fore_color.rgb = INK
            elif ri in hi_rows:
                cell.fill.fore_color.rgb = ACCENT_LT
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 else SURF
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if ci in num_cols else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            r.font.name = MONO if (ci in num_cols and ri > 0) else FONT
            r.font.size = Pt(size)
            r.font.bold = (ri == 0 and header) or ri in hi_rows
            r.font.color.rgb = WHITE if (ri == 0 and header) else INK
    return shp


def callout(slide_, x, y, w, h, tag, text, tone="accent", size=15):
    bar, bg, tc = {
        "accent":  (ACCENT, ACCENT_LT, ACCENT),
        "caution": (CAUTION, CAUTION_LT, CAUTION),
        "fail":    (FAIL, FAIL_LT, FAIL),
    }[tone]
    rect(slide_, x, y, w, h, bg)
    rect(slide_, x, y, 0.05, h, bar)
    box(slide_, x + 0.28, y + 0.16, w - 0.55, 0.24, tag.upper(), 10.5,
        bold=True, color=tc)
    box(slide_, x + 0.28, y + 0.48, w - 0.55, h - 0.62, text, size, color=INK)


def bignum(slide_, x, y, value, label, color=ACCENT, vsize=72, w=3.2):
    box(slide_, x, y, w, 1.1, value, vsize, bold=True, color=color, font=FONT)
    box(slide_, x, y + 1.08, w, 0.5, label, 13, color=MUTED)



# ================================================================ TITLE
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, WHITE)
rect(s, 0, 0, 0.32, H, ACCENT)
box(s, 1.5, 1.75, 10.5, 0.3, "PHASE 2 REVIEW  ·  SEMESTER 3  ·  SUMMER PROGRESS",
    12, bold=True, color=ACCENT)
box(s, 1.5, 2.25, 10.6, 1.5,
    "Domain-Adaptive Zero-Shot Anomaly Detection Using Vision-Language Models",
    38, bold=True, color=INK, spacing=1.06)
rect(s, 1.5, 3.95, 2.0, 0.04, ACCENT)
box(s, 1.5, 4.25, 10.0, 0.5,
    "Adapting a video anomaly detector to a new environment by writing a "
    "sentence — with no training, no target data, and every model frozen.",
    16, color=INK_2, spacing=1.3)
box(s, 1.5, 5.55, 6.0, 0.35, "Prakash Kumar Sarangi", 17, bold=True, color=INK)
box(s, 1.5, 5.92, 7.0, 0.6,
    "M.Tech Computer Science and Engineering\n"
    "Supervisors: Dr. Pranesh Das  ·  Dr. Raju Hazari", 12.5, color=MUTED,
    spacing=1.35)
box(s, 1.5, 6.75, 7.0, 0.3, "National Institute of Technology Calicut", 12.5,
    bold=True, color=ACCENT)
box(s, W - R - 2.2, 6.75, 2.2, 0.3, "18 August 2026", 12.5, color=MUTED,
    align=PP_ALIGN.RIGHT)

# ================================================================ 1 PROBLEM
s = slide("A detector is tied to the place it learned", "The problem")
bullets(s, [
    ("Standard approach.", "Show the system weeks of ordinary footage from one "
     "camera until it learns what normal looks like there."),
    ("The cost.", "Move it anywhere else and it fails — a new site means new "
     "footage collection and a full retraining cycle."),
], 1.95, width=CW, size=16.5)

rect(s, L, 3.15, 5.55, 2.5, SURF)
box(s, L + 0.35, 3.42, 4.9, 0.3, "SHOPPING MALL", 12, bold=True, color=ACCENT)
box(s, L + 0.35, 3.85, 4.9, 1.5,
    "A forklift moving through the aisles is\nan immediate alarm.",
    17, color=INK, spacing=1.35)

rect(s, L + 5.95, 3.15, 5.55, 2.5, SURF)
box(s, L + 6.3, 3.42, 4.9, 0.3, "FACTORY FLOOR", 12, bold=True, color=CAUTION)
box(s, L + 6.3, 3.85, 4.9, 1.5,
    "The identical forklift is completely\nroutine.",
    17, color=INK, spacing=1.35)

box(s, L, 5.95, CW, 0.5,
    "Same footage. Opposite answers. The picture did not change — the rule did.",
    17, bold=True, color=INK, align=PP_ALIGN.CENTER)

# ================================================================ 2 SHIFT
s = slide("The literature solves the adjacent problem", "Research gap  ·  1 of 2")
box(s, L, 1.95, CW, 0.4,
    "Domain adaptation distinguishes between kinds of difference between places.",
    16, color=INK_2)

table(s, [
    ["Type of shift", "What moves", "Example"],
    ["Covariate shift", "Appearance", "A stop sign in fog versus sunshine. Still a stop sign."],
    ["Concept shift", "The rule itself", "A bicycle on a road versus on a footpath. Identical image, opposite label."],
], L, 2.55, CW, col_w=[2.9, 2.2, 6.53], size=14, hi_rows=(2,), row_h=0.72)

callout(s, L, 4.85, CW, 1.35, "What the surveys say",
        "“Concept shift is usually not a common problem in popular object "
        "classification… this review mainly focuses on covariate shift.”   "
        "— Liu et al., 2022.   Singhal et al. (2023) list stable p(y|x) as the "
        "first condition under which domain adaptation is justified.",
        tone="caution", size=14.5)

box(s, L, 6.45, CW, 0.4,
    "Fair for object recognition — a cat is a cat everywhere. "
    "But anomaly detection is built on concept shift.",
    16, bold=True, color=INK)

# ================================================================ 3 GAP
s = slide("Why that matters for anomaly detection", "Research gap  ·  2 of 2")
box(s, L, 1.9, CW, 0.4,
    "In anomaly detection, “normal” is defined by the deployment context — "
    "not by the object. That is the definition of the task.",
    16.5, color=INK_2, spacing=1.3)

table(s, [
    ["Input", "Domain A", "Domain B"],
    ["A person running", "Park — normal", "Bank vault — anomalous"],
    ["A person lying down", "Beach — normal", "Factory aisle — anomalous"],
    ["A vehicle", "Road — normal", "Walkway — anomalous"],
], L, 2.7, CW, col_w=[3.6, 4.0, 4.03], size=14.5, row_h=0.5)

callout(s, L, 5.0, CW, 1.35, "The gap we target",
        "Methods that align appearance cannot help here — the appearance is "
        "already identical. Two domains can share p(x) exactly while the "
        "labelling function differs. A 2026 position paper argues the same "
        "premise independently; what is missing is a test of whether supplied "
        "context does any work.", size=14.5)

# ================================================================ 4 APPROACH
s = slide("DA-ZVAD: adapt by writing a sentence", "Proposed framework")
_here = os.path.dirname(os.path.abspath(__file__))
img = next((c for c in (
    os.path.join(_here, "dazvad_architecture.png"),
    os.path.join(_here, "..", "docs", "06_presentations", "dazvad_architecture.png"),
    os.path.join(_here, "..", "docs", "09_paper", "dazvad_architecture.png"),
) if os.path.isfile(c)), "")
if img:
    s.shapes.add_picture(img, Inches(L), Inches(1.95), width=Inches(7.5))
box(s, L + 7.85, 1.95, 3.68, 0.3, "EVERY MODEL FROZEN", 11.5, bold=True,
    color=ACCENT)
bullets(s, [
    ("M1", "Frozen CLIP scores each frame against “normal” and “abnormal” text."),
    ("M2", "Moving average over time — no parameters."),
    ("M3", "The scene description. The adaptation mechanism."),
    ("M4", "Frozen LLaVA explains each detected event."),
], 2.45, left=L + 7.85, width=3.68, size=13, gap=0.78)

callout(s, L + 7.85, 5.65, 3.68, 1.15, "Deployment",
        "Moving to a new site = editing one sentence. No target data. "
        "No gradients.", size=13)

# ================================================================ 5 FROZEN
s = slide("Why freezing everything is the point", "Method  ·  design rationale")
bullets(s, [
    ("The risk.", "If the system learned even a little from the new site, and "
     "performance improved, we could not say what caused it — the sentence, or "
     "the learning."),
    ("Our design.", "No parameter anywhere is allowed to change. Exactly one "
     "thing in the system can vary: the text."),
    ("The consequence.", "Any measured difference is attributable to the "
     "sentence. There is no other candidate."),
], 2.0, size=16.5, gap=0.95)

callout(s, L, 4.75, CW, 1.5, "Identifiability",
        "This is what makes the adaptation claim testable rather than asserted. "
        "You cannot run our central experiment on the competing systems — their "
        "text is learned on source data, or entangled in an LLM prior, or "
        "accompanied by a trained adapter.", size=16)

box(s, L, 6.5, CW, 0.4,
    "It also means the claim can be falsified — which is the next slide.",
    15, italic=True, color=MUTED)

# ================================================================ 6 PROTOCOL
s = slide("A protocol that can prove us wrong", "Evaluation design")
box(s, L, 1.95, CW, 0.4,
    "Run the identical pipeline four times. Change only the sentence.",
    16.5, color=INK_2)

table(s, [
    ["Condition", "Sentence supplied", "Purpose"],
    ["none", "M3 disabled", "Lower reference"],
    ["generic", "“a generic scene”", "Controls for merely having context"],
    ["matched", "Correct description of the scene", "The proposed operating condition"],
    ["mismatched", "Description of the wrong domain", "FALSIFYING CONTROL"],
], L, 2.6, CW, col_w=[2.4, 4.6, 4.63], size=14, hi_rows=(4,), row_h=0.52)

callout(s, L, 5.15, CW, 1.35, "Predicted signature, fixed before measurement",
        "matched ≥ generic ≥ none, with mismatched measurably WORSE. "
        "If a deliberately wrong description costs nothing, the method ignores "
        "its context and our claim is refuted.", size=15.5)

# ================================================================ 7 SETUP
s = slide("What we ran", "Experiments  ·  setup")
table(s, [
    ["Benchmarks", "ShanghaiTech (13 views) and CUHK Avenue (1 view)"],
    ["Scale", "128 test clips · 28,118 frames · frame-level ground truth"],
    ["Hardware", "NVIDIA A40, college GPU server"],
    ["Backbone", "CLIP ViT-L/14 (LAION-2B), frozen · PyTorch 2.3.1 / CUDA 12.1"],
    ["Runs", "5 full experiment runs + cached-embedding analysis"],
], L, 2.0, CW, col_w=[2.4, 9.23], size=15, header=False, row_h=0.52)

callout(s, L, 4.8, CW, 1.75, "Every run records its own conditions",
        "Each run writes a manifest: the exact code commit, whether the working "
        "tree was clean, host, GPU, driver and library versions, full "
        "configuration, and per-dataset frame and label counts. Committed "
        "alongside the results — every figure in this deck is traceable to the "
        "state that produced it.", size=15)

# ================================================================ 8 FAILURE
s = slide("The first result was a failure", "Experiments  ·  what happened")
rect(s, L, 1.95, CW, 1.9, FAIL_LT)
rect(s, L, 1.95, 0.05, 1.9, FAIL)
bignum(s, L + 0.5, 2.15, "0.49", "frame-level AUROC — 0.50 is random guessing",
       color=FAIL, vsize=64, w=3.0)
box(s, L + 4.2, 2.35, 7.2, 1.2,
    "Fifty minutes of GPU time produced a detector performing exactly as well "
    "as a coin flip.", 17, color=INK, spacing=1.3)

box(s, L, 4.2, CW, 0.4, "And the central experiment came out backwards:",
    16.5, bold=True, color=INK)

table(s, [
    ["Condition", "Predicted", "Observed"],
    ["matched (correct description)", "Best", "WORST — 0.666"],
    ["mismatched (wrong description)", "Worst", "Among the best — 0.695"],
], L, 4.8, CW, col_w=[5.03, 3.3, 3.3], size=14.5, row_h=0.5)

box(s, L, 6.55, CW, 0.4,
    "This is the point at which the idea looks broken.",
    16, italic=True, color=FAIL)

# ================================================================ 9 FIX 1
s = slide("Diagnosis 1 — we were measuring it wrong", "Experiments  ·  diagnosis")
bullets(s, [
    ("13 camera views.", "CLIP sits at a different baseline score under each one "
     "— different lighting, different angle."),
    ("Our error.", "We pooled every frame from all 13 cameras into one ranking."),
    ("The analogy.", "Ranking students from different schools by raw marks when "
     "the schools grade differently. The comparison destroys the ordering."),
], 1.95, size=16, gap=0.88)

box(s, L, 4.5, CW, 0.35,
    "The published protocol for this benchmark normalises each clip first. "
    "We were not doing it.", 15.5, color=MUTED)

rect(s, L, 5.05, CW, 1.35, ACCENT_LT)
bignum(s, L + 0.6, 5.12, "0.49", "as reported", color=MUTED, vsize=42, w=2.2)
box(s, L + 3.0, 5.42, 0.9, 0.6, "→", 34, bold=True, color=ACCENT)
bignum(s, L + 4.2, 5.12, "0.71", "correctly pooled", color=ACCENT, vsize=42, w=2.4)
box(s, L + 7.3, 5.35, 4.1, 0.9,
    "Identical scores. Uses no labels.\nWe report both figures in the paper.",
    14, color=INK, spacing=1.3)

# ================================================================ 10 FIX 2
s = slide("Diagnosis 2 — the description cancelled itself out",
          "Experiments  ·  the finding")
box(s, L, 1.9, CW, 0.35,
    "We were appending the scene sentence to BOTH prompt sets:", 16, color=INK_2)

rect(s, L, 2.45, 5.55, 1.15, SURF)
box(s, L + 0.3, 2.62, 5.0, 0.25, "NORMAL PROMPT", 10.5, bold=True, color=ACCENT)
box(s, L + 0.3, 2.95, 5.0, 0.6,
    "“a campus walkway with pedestrians,\n everything is normal”",
    13, color=INK, font=MONO, spacing=1.3)

rect(s, L + 5.95, 2.45, 5.55, 1.15, SURF)
box(s, L + 6.25, 2.62, 5.0, 0.25, "ABNORMAL PROMPT", 10.5, bold=True, color=CAUTION)
box(s, L + 6.25, 2.95, 5.0, 0.6,
    "“a campus walkway with pedestrians,\n but something is wrong”",
    13, color=INK, font=MONO, spacing=1.3)

bullets(s, [
    ("Shared words.", "Each prompt set is averaged into one summary vector. "
     "Shared text enters both — so the two summaries move toward each other."),
    ("The method depends on them being different.", "We were erasing the very "
     "contrast the decision rests on."),
    ("Which explains the inversion.", "An accurate description matches every "
     "frame strongly, so it absorbs the most contrast. A wrong one matches "
     "nothing, so it does no damage."),
], 3.85, size=15.5, gap=0.78)

callout(s, L, 6.25, CW, 0.85, "The fix",
        "Attach the description to the normal prompts only — the scene defines "
        "what normal looks like here; an anomaly is a departure from it.",
        size=15)

# ================================================================ 11 SWEEP
s = slide("After the fix: the predicted signature", "Results  ·  central experiment")
box(s, L, 1.9, CW, 0.35,
    "ShanghaiTech, all 107 clips, per-clip normalised. Every model frozen; "
    "only the sentence and its injection point vary.", 14.5, color=MUTED)

table(s, [
    ["Injection point", "none", "generic", "matched", "mismatched", "gap"],
    ["Both prompt sets", "0.707", "0.670", "0.666", "0.695", "−0.029"],
    ["Normal set only", "0.707", "0.691", "0.734", "0.628", "+0.105"],
], L, 2.5, CW, col_w=[3.23, 1.68, 1.68, 1.68, 1.68, 1.68], size=15,
    hi_rows=(2,), row_h=0.58, num_cols=(1, 2, 3, 4, 5))

box(s, L, 4.35, CW, 0.35,
    "The “none” column is identical in both rows — confirming nothing but the "
    "injection point changed.", 14, italic=True, color=MUTED)

rect(s, L, 4.95, 5.55, 1.75, ACCENT_LT)
bignum(s, L + 0.45, 5.1, "0.734", "correct description", color=ACCENT,
       vsize=44, w=3.5)

rect(s, L + 5.95, 4.95, 5.55, 1.75, FAIL_LT)
bignum(s, L + 6.4, 5.1, "0.628", "wrong description  —  a 10-point penalty",
       color=FAIL, vsize=44, w=4.9)

box(s, L, 6.9, CW, 0.35,
    "Nothing else in the system was permitted to change. The text caused it.",
    15.5, bold=True, color=INK, align=PP_ALIGN.CENTER)

# ================================================================ 12 PRECISION
s = slide("Stating the claim precisely", "Results  ·  interpretation")
rect(s, L, 1.95, 5.55, 2.15, SURF)
box(s, L + 0.4, 2.2, 4.8, 0.3, "THE WEAKER HALF", 11.5, bold=True, color=MUTED)
box(s, L + 0.4, 2.6, 4.8, 1.3,
    "A correct description beats none by +0.027.\n\n"
    "Positive at every window, but inside the ±0.036 split-to-split spread.",
    15, color=INK, spacing=1.3)

rect(s, L + 5.95, 1.95, 5.55, 2.15, ACCENT_LT)
box(s, L + 6.35, 2.2, 4.8, 0.3, "THE CLAIM WE MAKE", 11.5, bold=True, color=ACCENT)
box(s, L + 6.35, 2.6, 4.8, 1.3,
    "A WRONG description costs −0.105.\n\n"
    "Unambiguous, and nothing else could have caused it.",
    15, color=INK, spacing=1.3)

callout(s, L, 4.4, CW, 1.4, "How to read the mechanism",
        "The description constrains a decision boundary rather than adding "
        "information. It does not reliably lift performance when correct; it "
        "degrades performance sharply when misdirected.", size=16)

box(s, L, 6.15, CW, 0.7,
    "A secondary finding, not present in the literature: WHERE the description "
    "is injected dominates WHAT it says — to the point of reversing the "
    "direction of the effect.",
    16, bold=True, color=INK, spacing=1.3)

# ================================================================ 13 ABLATION
s = slide("Which components earn their place", "Results  ·  component ablation")
table(s, [
    ["Scoring signal", "Held-out AUROC", "Full test set"],
    ["Scene-centre normality only", "0.585 ± 0.025", "0.585"],
    ["Semantic + scene-centre", "0.645 ± 0.034", "0.640"],
    ["Kinematic — motion only", "0.685 ± 0.015", "0.686"],
    ["Semantic + kinematic", "0.711 ± 0.034", "0.706"],
    ["Semantic — language only", "0.718 ± 0.036", "0.707"],
], L, 2.0, CW, col_w=[5.2, 3.25, 3.18], size=14.5, hi_rows=(5,), row_h=0.47,
    num_cols=(1, 2))

box(s, L, 4.65, CW, 0.35,
    "Held-out figures are means over five clip-level partitions, with the "
    "spread across them.", 13.5, italic=True, color=MUTED)

bullets(s, [
    ("The language pathway alone is best.", "Adding a motion signal costs 0.001 "
     "on the full set; adding scene-centre normality costs 0.067."),
    ("This was not the expected answer.", "ShanghaiTech's anomalies look "
     "kinematic — a bicycle at cycling speed — so appearance and motion should "
     "be complementary. Measured here, they are not."),
], 5.2, size=15, gap=1.0)

# ================================================================ 14 AVENUE
s = slide("A second domain — the replication test", "Results  ·  cross-domain")
box(s, L, 1.9, CW, 0.35,
    "Identical frozen configuration applied to CUHK Avenue. Nothing retuned. "
    "Only the scene sentence changed.", 15, color=INK_2)

table(s, [
    ["Dataset", "none", "generic", "matched", "mismatched", "gap"],
    ["ShanghaiTech  (13 views)", "0.707", "0.691", "0.734", "0.628", "+0.105"],
    ["CUHK Avenue  (1 view)", "0.706", "0.729", "0.677", "0.657", "+0.020"],
], L, 2.5, CW, col_w=[3.23, 1.68, 1.68, 1.68, 1.68, 1.68], size=15,
    hi_rows=(2,), row_h=0.58, num_cols=(1, 2, 3, 4, 5))

rect(s, L, 4.35, 5.55, 1.5, ACCENT_LT)
box(s, L + 0.35, 4.55, 4.9, 0.3, "DETECTION TRANSFERS", 11.5, bold=True, color=ACCENT)
box(s, L + 0.35, 4.9, 4.9, 0.8,
    "0.706 vs 0.707 with no descriptor.\nThe detector works equally well.",
    15, color=INK, spacing=1.3)

rect(s, L + 5.95, 4.35, 5.55, 1.5, FAIL_LT)
box(s, L + 6.3, 4.55, 4.9, 0.3, "ADAPTATION DOES NOT", 11.5, bold=True, color=FAIL)
box(s, L + 6.3, 4.9, 4.9, 0.8,
    "Gap 5× smaller, and the correct\ndescription scores BELOW none.",
    15, color=INK, spacing=1.3)

callout(s, L, 6.0, CW, 1.15, "Our explanation — stated as a conjecture",
        "ShanghaiTech has 13 camera views; Avenue has one. A scene description "
        "has work to do only when there are several environments to tell "
        "apart. Testable: run the sweep inside a single ShanghaiTech view.",
        tone="caution", size=14.5)

# ================================================================ 15 WITHIN-VIEW
s = slide("We tested that explanation, and it held", "Results  ·  the control")
box(s, L, 1.9, CW, 0.7,
    "ShanghaiTech is thirteen single-view datasets stacked together. If the "
    "descriptor works by telling the model WHICH scene it is in, then confining "
    "the sweep to one camera view should reproduce Avenue's flat result.",
    15.5, color=INK_2, spacing=1.3)

table(s, [
    ["Evaluation", "Views", "Clips", "Gap"],
    ["Pooled across views", "13", "107", "+0.105"],
    ["Within a single view (mean)", "9", "5–34 each", "+0.034"],
    ["CUHK Avenue (single view)", "1", "21", "+0.020"],
], L, 2.85, CW, col_w=[5.6, 1.9, 2.3, 1.83], size=15, hi_rows=(2,), row_h=0.56,
    num_cols=(3,))

callout(s, L, 4.9, CW, 1.15, "The prediction could have failed",
        "A within-view gap near +0.105 would have refuted the explanation "
        "outright. It came back at a third of that, next to Avenue's figure.",
        size=15.5)

box(s, L, 6.25, CW, 0.7,
    "What the sentence mainly supplies is WHICH scene you are in — not what "
    "counts as normal within it. Three of nine views show a negative gap, so "
    "inside one scene the effect is not reliable.",
    15.5, bold=True, color=INK, spacing=1.3)

# ================================================================ 16 NEGATIVES
s = slide("Six things that did not work", "Results  ·  negative findings")
table(s, [
    ["Modification", "Outcome"],
    ["Quadrant scoring, to catch small objects", "No improvement in any configuration"],
    ["Prompts naming bicycles and vehicles", "Much worse alone — 0.486"],
    ["Clip's own average as the normality reference", "0.585, and it degrades the language signal"],
    ["Adding a motion signal", "Costs 0.001 — not complementary"],
    ["Local temporal deviation before projection", "Better scorer, but narrows the context gap"],
    ["Per-prompt max pooling", "0.678 against 0.707"],
], L, 2.0, CW, col_w=[6.6, 5.03], size=13.5, row_h=0.5)

callout(s, L, 5.3, CW, 1.55, "Why these are in the deck",
        "The claim is that the minimal configuration is the right one. That is "
        "only credible alongside the alternatives that were tried. A clean "
        "table of successes is the artefact that is easy to fabricate; six "
        "diagnosed failures are not.", size=15.5)

# ================================================================ 16 POSITION
s = slide("Where this sits against the literature", "Assessment")
table(s, [
    ["Method", "AUROC", "Cost per frame"],
    ["Liu et al. 2018 — trained on the target scene", "≈ 0.728", "1 trained model"],
    ["LAVAD (CVPR 2024), training-free", "≈ 0.85", "Captioner + LLM + refiner"],
    ["DA-ZVAD (ours), training-free", "0.734", "One frozen encoder + a sentence"],
], L, 2.0, CW, col_w=[5.6, 2.3, 3.73], size=14.5, hi_rows=(3,), row_h=0.56)

bullets(s, [
    ("The fair comparison is the trained baseline.", "We match the benchmark's "
     "own 2018 baseline while using none of its training data."),
    ("We do not match LAVAD, and say so.", "It runs three large models per "
     "frame. We run in about 7 GB."),
], 4.2, size=16, gap=1.0)

callout(s, L, 6.0, CW, 1.15, "Honest position",
        "We do not expect to exceed trained state-of-the-art detectors on "
        "absolute AUROC, and the paper makes no such claim.", size=16)

# ================================================================ 17 LIMITS
s = slide("Limitations we are stating ourselves", "Assessment")
bullets(s, [
    ("The mechanism is bounded.", "It works on ShanghaiTech and nearly vanishes "
     "on Avenue. The scene-diversity explanation rests on two datasets and is "
     "not established."),
    ("Resolution ceiling.", "Whole-frame embeddings at 224×224 cannot resolve "
     "small objects. Quadrant scoring did not close it; patch-level scoring is "
     "untested."),
    ("The metric is scale-sensitive.", "Per-clip normalisation is affine, so a "
     "monotone rescaling of the score changes the pooled figure even though the "
     "ranking is identical. We found this the hard way."),
    ("Configuration selection.", "No validation split exists for these "
     "benchmarks, so we split clips and report the half never used to select."),
    ("Both video domains are outdoor pedestrian surveillance.", "The "
     "industrial contrast is evaluated for detection but not for the sweep."),
    ("M4 not yet run.", "The explanation module has produced no video results."),
], 1.95, size=14.5, gap=0.86)

# ================================================================ 18 NEXT
s = slide("Phase 3 plan", "Next")
table(s, [
    ["Priority", "Work", "Why it matters"],
    ["1", "Sweep within a single ShanghaiTech camera view",
     "Settles whether scene diversity is the operative variable"],
    ["2", "Context sweep on MVTec AD",
     "Gives the industrial–surveillance contrast for adaptation, not just detection"],
    ["3", "Patch-level scoring against spatial tokens",
     "Most likely route to a materially higher figure"],
    ["4", "Run M4 — explanations under matched vs mismatched context",
     "Completes the framework; supplies the qualitative result"],
], L, 2.0, CW, col_w=[1.2, 5.5, 4.93], size=13.5, row_h=0.62, num_cols=(0,))

callout(s, L, 5.2, CW, 1.5, "Infrastructure is now built",
        "Server, datasets, driver and embedding cache are in place. Encoding a "
        "benchmark once takes 21 minutes, after which a new scoring hypothesis "
        "is evaluated in seconds rather than 50 minutes.", size=15.5)

# ================================================================ 19 SUMMARY
s = slide("Summary", "Phase 2")
items = [
    ("The gap", "Domain adaptation research targets covariate shift by explicit "
     "scoping. Anomaly detection is dominated by concept shift."),
    ("The method", "Every model frozen; adaptation carried entirely by a "
     "sentence — which is what makes the claim identifiable."),
    ("The result", "0.734 AUROC on ShanghaiTech with no training, and a wrong "
     "description costs 0.105 — evidence the text is load-bearing."),
    ("The finding", "Where the description is injected dominates what it says, "
     "to the point of reversing the effect. Not in the literature."),
    ("The boundary", "It does not replicate on a single-scene benchmark, and a "
     "within-view control shows why: the sentence tells the model which scene "
     "it is in. The claim is scoped, and the scoping is measured."),
]
y = 1.95
for i, (head, body) in enumerate(items, 1):
    rect(s, L, y, CW, 0.86, SURF if i % 2 else WHITE)
    box(s, L + 0.3, y + 0.24, 0.5, 0.4, str(i), 22, bold=True, color=ACCENT)
    box(s, L + 0.95, y + 0.12, 2.5, 0.3, head, 14, bold=True, color=INK)
    box(s, L + 0.95, y + 0.42, 10.4, 0.4, body, 13.5, color=INK_2, spacing=1.2)
    y += 0.9

box(s, L, 6.75, CW, 0.35,
    "All results reproducible from committed code and run manifests.",
    13, italic=True, color=MUTED, align=PP_ALIGN.CENTER)

# ================================================================ NOTES
NOTES = [
 "Title. Introduce yourself and the one-line premise: adapting an anomaly "
 "detector to a new place by writing a sentence, with nothing retrained.",

 "Open with the example, not the definition. Mall versus factory - the same "
 "forklift, opposite answers. Then the cost: every new customer means new "
 "footage and a retraining cycle. That cost is what we are removing.",

 "This is the slide that shows you read the six surveys. Read the Liu quote "
 "aloud - the field scopes concept shift OUT by explicit decision. Make clear "
 "you are not criticising them: for object recognition it is the right call.",

 "Land the consequence. In anomaly detection normality IS the deployment "
 "context. Mention the 2026 position paper yourself - an independent group "
 "reached the same premise, which answers 'did you invent this problem?'",

 "Walk the four modules left to right in about thirty seconds. Then stop on "
 "M3 - the orange box - and say that is the only place the domain enters.",

 "The strongest methodological point in the deck. Say it slowly. And note that "
 "the competing systems CANNOT run our experiment - their text is learned on "
 "source data, or entangled in an LLM prior, or paired with a trained adapter.",

 "Emphasise that the interpretation was fixed BEFORE any measurement, and that "
 "the mismatched condition is designed to refute us. A panel will respect a "
 "test you could have failed - and later in the deck, one you partly did.",

 "Keep this brief - it is the credibility slide. If anyone doubts the work is "
 "yours, offer to open a manifest. Five runs, all recorded.",

 "Do not rush past this and do not apologise for it. The first run was chance, "
 "and the key experiment came out backwards. Then pause. The next two slides "
 "are what makes this presentation worth listening to.",

 "The schools analogy works on everyone - use it. Stress twice that the fix "
 "uses no labels and is the benchmark's own published protocol, and that both "
 "figures appear in the paper.",

 "The most interesting slide. Point at the two prompts and let them see the "
 "shared words. Then the punchline: the more accurate the description, the "
 "more damage it does. That is why the result inverted.",

 "The headline. Point at the none column being identical in both rows - that "
 "is the control, and it proves only the injection point changed. Then point "
 "at 0.628 and say a wrong sentence costs ten points.",

 "Be scrupulous. The +0.027 is positive at every window but sits inside the "
 "spread, so we report direction and not magnitude. The claim rests on the "
 "mismatched penalty. Understating here protects you.",

 "Volunteer the surprise: we expected motion to help and it does not. The "
 "pooled embedding already registers enough of it. Reporting the prediction "
 "that failed is stronger than reporting only the ones that held.",

 "The most important slide in the deck. Detection transfers almost exactly; "
 "the adaptation does not. Give the honest reading - on Avenue a placeholder "
 "beats an accurate description, so that benefit cannot be domain adaptation. "
 "Then give the conjecture AND the experiment that would settle it.",

 "This is the slide that shows a full cycle: we saw something odd, formed an "
 "explanation, designed a test that could have killed it, ran it, and it held. "
 "Say out loud that a gap near +0.105 here would have refuted us. Then give "
 "the caveat yourself - three of nine views are negative, so within one scene "
 "the effect is not reliable.",

 "Say why you are showing failures: the claim is that the minimal "
 "configuration is right, and that is only credible next to what was tried. "
 "If the panel suspects generated work, this slide is your best answer.",

 "Anchor on the trained baseline, not on LAVAD. You match the benchmark's own "
 "2018 baseline using none of its training data. Then concede LAVAD openly.",

 "Deliver these confidently rather than apologetically. The metric limitation "
 "is worth dwelling on - it is a real methodological point about a protocol "
 "the whole field uses, and we found it by chasing our own bug.",

 "Priority one is the within-view control because it turns our explanation "
 "from a conjecture into a result. Note the infrastructure cost is now paid.",

 "Close on the boundary, not the number. The claim is scoped rather than "
 "abandoned, and we named the experiment that settles it. Then hand over.",
]
for sl, txt in zip(prs.slides, NOTES):
    sl.notes_slide.notes_text_frame.text = txt

# ================================================================ SAVE
out = os.path.join(_here, "DA-ZVAD_Phase2_Review.pptx")
prs.save(out)
print(f"saved: {out}")
print(f"slides: {len(prs.slides._sldIdLst)}   notes: {len(NOTES)}")
