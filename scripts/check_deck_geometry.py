# -*- coding: utf-8 -*-
"""Geometry check for the Phase 2 deck.

There is no renderer on this machine, so overflow cannot be seen. This measures
instead: every shape against the slide bounds, and every text box against the
real width of its own text, wrapped with the actual font metrics. It reports
where text needs more vertical room than its box was given, which is the
failure that silently produces overlapping lines in PowerPoint.
"""
import sys
from pptx import Presentation
from pptx.util import Emu
from PIL import ImageFont

EMU = 914400.0
SEGOE = r"C:\Windows\Fonts\segoeui.ttf"
SEGOE_B = r"C:\Windows\Fonts\segoeuib.ttf"
_cache = {}


def font(size_pt, bold):
    key = (round(size_pt, 1), bold)
    if key not in _cache:
        # 1pt = 1/72in; PIL sizes in px, so render at 72dpi to make px == pt.
        _cache[key] = ImageFont.truetype(SEGOE_B if bold else SEGOE,
                                         max(1, int(round(size_pt))))
    return _cache[key]


def text_width_in(txt, size_pt, bold):
    f = font(size_pt, bold)
    return f.getlength(txt) / 72.0


def wrapped_lines(runs, width_in):
    """Greedy wrap of a paragraph's runs into width_in inches."""
    words, sizes = [], []
    for txt, size, bold in runs:
        for w in txt.split(" "):
            if w:
                words.append(w)
                sizes.append((size, bold))
    if not words:
        return 0, 0.0
    lines, cur, cur_w, max_size = 1, "", 0.0, 0.0
    biggest = max(s for s, _ in sizes)
    for w, (size, bold) in zip(words, sizes):
        piece = (w if not cur else " " + w)
        pw = text_width_in(piece, size, bold)
        if cur and cur_w + pw > width_in:
            lines += 1
            cur, cur_w = w, text_width_in(w, size, bold)
        else:
            cur += piece
            cur_w += pw
    return lines, biggest


def check(path):
    prs = Presentation(path)
    SW, SH = prs.slide_width / EMU, prs.slide_height / EMU
    problems = []
    for i, slide in enumerate(prs.slides, 1):
        for shp in slide.shapes:
            try:
                x, y = shp.left / EMU, shp.top / EMU
                w, h = shp.width / EMU, shp.height / EMU
            except TypeError:
                continue
            if x < -0.02 or y < -0.02 or x + w > SW + 0.02 or y + h > SH + 0.02:
                problems.append(
                    f"slide {i:>2}: {shp.shape_type} out of bounds  "
                    f"x={x:.2f} y={y:.2f} w={w:.2f} h={h:.2f}")
            if not shp.has_text_frame:
                continue
            for p in shp.text_frame.paragraphs:
                runs = [(r.text, (r.font.size.pt if r.font.size else 18),
                         bool(r.font.bold)) for r in p.runs if r.text]
                if not runs:
                    continue
                nlines, size = wrapped_lines(runs, w)
                spacing = p.line_spacing if isinstance(p.line_spacing, float) else 1.0
                need = nlines * size * spacing * 1.22 / 72.0
                if need > h + 0.02:
                    txt = "".join(t for t, _, _ in runs)
                    problems.append(
                        f"slide {i:>2}: text needs {need:.2f}in in a {h:.2f}in "
                        f"box ({nlines} lines @ {size:.1f}pt) "
                        f"y={y:.2f} -> bottom {y + need:.2f}  |  {txt[:58]!r}")
    return prs, SW, SH, problems


if __name__ == "__main__":
    prs, SW, SH, problems = check(sys.argv[1])
    print(f"slide size {SW:.2f} x {SH:.2f} in,  {len(prs.slides._sldIdLst)} slides")
    if not problems:
        print("no overflow detected")
    else:
        print(f"{len(problems)} issue(s):")
        for p in problems:
            print("  " + p)
