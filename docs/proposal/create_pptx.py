import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# === REFINED MODERN COLOR PALETTE ===
BG_BASE = RGBColor(0xFF, 0xFF, 0xFF)
BG_CARD = RGBColor(0xF1, 0xF5, 0xF9)
ACCENT_INDIGO = RGBColor(0x43, 0x38, 0xCA)
ACCENT_TEAL = RGBColor(0x0D, 0x94, 0x88)
ACCENT_ROSE = RGBColor(0xE1, 0x1D, 0x48)
ACCENT_AMBER = RGBColor(0xD9, 0x77, 0x06)
ACCENT_VIOLET = RGBColor(0x7C, 0x3A, 0xED)
TEXT_MAIN = RGBColor(0x0F, 0x17, 0x2A)
TEXT_MUTED = RGBColor(0x47, 0x55, 0x69)
TEXT_DIM = RGBColor(0x94, 0xA3, 0xB8)

FONT_NAME = 'Segoe UI'
FONT_TITLE = 'Segoe UI Semibold'

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color=BG_BASE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, font_size=18, color=TEXT_MAIN, bold=False, align=PP_ALIGN.LEFT, font=FONT_NAME):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return txBox

def add_bullets(slide, left, top, width, height, items, font_size=16, color=TEXT_MUTED, font=FONT_NAME):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font
        p.space_after = Pt(10)
    return txBox

def add_rect(slide, left, top, width, height, color=BG_CARD, alpha=1.0, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Adjust border radius if it's a rounded rect
    if shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        shape.adjustments[0] = 0.05
    return shape

def add_line(slide, left, top, width, color=ACCENT_INDIGO):
    add_rect(slide, left, top, width, 0.04, color, shape_type=MSO_SHAPE.RECTANGLE)

def create_header(slide, title, subtitle=None, color=ACCENT_INDIGO):
    add_text(slide, 0.8, 0.4, 10, 0.7, title, font_size=32, color=TEXT_MAIN, bold=True, font=FONT_TITLE)
    add_line(slide, 0.8, 1.1, 1.5, color)
    if subtitle:
        add_text(slide, 0.8, 1.25, 10, 0.4, subtitle, font_size=16, color=color, bold=True)

# ===== SLIDE 1: TITLE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
# Decorative tech nodes in background
for i in range(5):
    add_rect(s, 0.5+i*2.5, 7.3, 2.0, 0.05, TEXT_DIM, shape_type=MSO_SHAPE.RECTANGLE)

add_line(s, 2, 2.0, 9.333, ACCENT_INDIGO)
add_text(s, 1, 2.3, 11.333, 1.5, "Multimodal Anomaly Detection in Industrial\nQuality Inspection using Vision-Language Models", font_size=34, color=TEXT_MAIN, bold=True, align=PP_ALIGN.CENTER, font=FONT_TITLE)
add_text(s, 1, 3.8, 11.333, 0.6, "Literature Survey & Proposed Research Direction", font_size=20, color=ACCENT_TEAL, align=PP_ALIGN.CENTER, bold=True)

add_text(s, 1, 5.2, 11.333, 0.5, "[Your Name]  •  [Roll Number]", font_size=16, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
add_text(s, 1, 5.6, 11.333, 0.5, "Guide: Prof. [Professor Name]", font_size=15, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
add_text(s, 1, 6.0, 11.333, 0.5, "Department of Computer Science & Engineering | NIT Calicut", font_size=14, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# ===== SLIDE 2: OUTLINE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
create_header(s, "Presentation Outline", None, ACCENT_INDIGO)

add_rect(s, 1.5, 1.8, 10.333, 4.8, BG_CARD)
add_bullets(s, 2.0, 2.2, 9.333, 4, [
    "Introduction: Problem Context & Motivation",
    "Background Concepts (CLIP, MLLMs, MVTec AD)",
    "Literature Survey: Review of 5 Core Papers",
    "Comparative Analysis & Synthesized Insights",
    "Identified Research Gaps",
    "Proposed Research Methodology & Pipeline",
    "Expected Contributions & Project Timeline"
], font_size=20, color=TEXT_MAIN)

# Add numeric badges
for i in range(7):
    y = 2.2 + i * 0.585
    r = add_rect(s, 1.3, y - 0.05, 0.4, 0.4, ACCENT_INDIGO, shape_type=MSO_SHAPE.OVAL)
    add_text(s, 1.3, y - 0.05, 0.4, 0.4, str(i+1), font_size=14, color=TEXT_MAIN, bold=True, align=PP_ALIGN.CENTER)

# ===== SLIDE 3: PROBLEM =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
create_header(s, "Problem & Motivation", "Industrial Quality Inspection Needs an Upgrade", ACCENT_TEAL)

add_rect(s, 0.8, 2.0, 5.5, 2.5, BG_CARD)
add_text(s, 1.1, 2.2, 5.1, 0.4, "Current Limitations", font_size=18, color=ACCENT_ROSE, bold=True)
add_bullets(s, 1.1, 2.7, 5.1, 1.5, [
    "Single-modality approaches require extensive training data",
    "Models must be retrained per product category",
    "Lacks interpretability (no natural language explanations)",
    "Human inspection remains slow & error-prone"
], font_size=15)

add_rect(s, 7.0, 2.0, 5.5, 2.5, BG_CARD)
add_text(s, 7.3, 2.2, 5.1, 0.4, "The VLM Opportunity", font_size=18, color=ACCENT_TEAL, bold=True)
add_bullets(s, 7.3, 2.7, 5.1, 1.5, [
    "Vision-Language Models enable zero-shot anomaly detection",
    "MLLMs provide human-like explanations and reasoning",
    "No product-specific training data needed (scalable)",
    "Enables intelligent, natural interactions with AI inspectors"
], font_size=15)

# Research question box
add_rect(s, 1.5, 5.0, 10.333, 1.5, ACCENT_VIOLET)
add_text(s, 1.7, 5.2, 10.0, 0.4, "Core Research Question", font_size=16, color=TEXT_MAIN, bold=True)
add_text(s, 1.7, 5.6, 9.9, 0.8, "How can VLMs and MLLMs be leveraged for zero-shot anomaly detection in industrial inspection, achieving robust detection AND interpretable explanations without task-specific training?", font_size=17, color=TEXT_MAIN, align=PP_ALIGN.CENTER, bold=True)

# ===== SLIDE 4: BACKGROUND =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
create_header(s, "Technical Background", "Foundational Concepts & Datasets", ACCENT_INDIGO)

for (title, items_list, x, clr) in [
    ("CLIP (OpenAI, 2021)", ["Trained on 400M image-text pairs", "Contrastive learning aligns vision/text", "Shared embedding space", "Enables zero-shot matching"], 0.8, ACCENT_INDIGO),
    ("MLLMs (e.g., LLaVA)", ["Vision encoder combined with LLM", "Processes visual inputs & texts", "Can describe, explain, and reason", "Emerging open-source capabilities"], 4.8, ACCENT_TEAL),
    ("MVTec AD Benchmark", ["Gold standard for industrial AD", "15 manufacturing categories", "High-res images with pixel labels", "Metric: Area Under ROC (AUROC)"], 8.8, ACCENT_AMBER),
]:
    add_rect(s, x, 2.0, 3.7, 3.5, BG_CARD)
    add_rect(s, x, 2.0, 3.7, 0.6, clr, shape_type=MSO_SHAPE.RECTANGLE)
    add_text(s, x+0.2, 2.1, 3.3, 0.4, title, font_size=18, color=TEXT_MAIN, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(s, x+0.2, 2.8, 3.3, 2.5, items_list, font_size=15)

# ===== Helper for Paper Slides =====
def create_paper_slide(prs, title, paper_meta, method_points, results_text, limitations, takeaway, color=ACCENT_INDIGO):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s)
    add_text(s, 0.8, 0.4, 11, 0.4, paper_meta, font_size=16, color=color, bold=True)
    add_text(s, 0.8, 0.8, 11, 0.6, title, font_size=28, color=TEXT_MAIN, bold=True, font=FONT_TITLE)
    add_line(s, 0.8, 1.4, 11.7, color)

    # Method
    add_rect(s, 0.8, 1.8, 7.5, 3.6, BG_CARD)
    add_text(s, 1.1, 2.0, 7.1, 0.4, "Methodology", font_size=18, color=color, bold=True)
    add_bullets(s, 1.1, 2.5, 6.9, 2.7, method_points, font_size=15)

    # Results & Limitations
    add_rect(s, 8.5, 1.8, 4.0, 1.5, BG_CARD)
    add_text(s, 8.7, 2.0, 3.6, 0.3, "Key Results", font_size=16, color=ACCENT_TEAL, bold=True)
    add_text(s, 8.7, 2.4, 3.6, 0.8, results_text, font_size=16, color=TEXT_MAIN, bold=True)

    add_rect(s, 8.5, 3.5, 4.0, 1.9, BG_CARD)
    add_text(s, 8.7, 3.7, 3.6, 0.3, "Limitations", font_size=16, color=ACCENT_ROSE, bold=True)
    add_bullets(s, 8.7, 4.1, 3.6, 1.2, limitations, font_size=14)

    # Takeaway Footer
    add_rect(s, 0.8, 5.8, 11.7, 1.0, clr)
    add_rect(s, 0.82, 5.82, 11.66, 0.96, BG_BASE) # Inner cut
    add_text(s, 1.1, 6.1, 11.1, 0.5, f"Takeaway: {takeaway}", font_size=16, color=TEXT_MAIN, bold=True)
    return s

# Papers
create_paper_slide(prs, "PatchCore: Towards Total Recall in Industrial AD", "Paper 1 | Roth et al. [1] | CVPR 2022",
    ["Extracts patch features from ImageNet-pretrained CNN", "Stores diverse normal features in a memory bank (coreset subsampling)", "Anomalies detected via K-Nearest Neighbor (KNN) distance", "State-of-the-art traditional upper-bound baseline"],
    "99.1% Image AUROC \nStrongest generic baseline on MVTec.",
    ["Requires training per category", "No natural language explanations", "Cannot generalize zero-shot"],
    "PatchCore sets the accuracy ceiling, but highlights the need for training-free, text-aware methods.", ACCENT_INDIGO)

create_paper_slide(prs, "WinCLIP: Zero-/Few-Shot Anomaly Classification & Segmentation", "Paper 2 | Jeong et al. [2] | CVPR 2023",
    ["First zero-shot application of CLIP to industrial inspection", "Compositional Prompt Ensembles (CPE): e.g., 'a damaged [object]'", "Multi-scale sliding window extracts patch-level CLIP features", "Compares text vs visual features to locate anomalies"],
    "91.8% Image AUROC \n85.1% Pixel AUROC (zero-shot)",
    ["Relies on ad-hoc, hand-crafted text prompts", "Lacks reasoning & explanations"],
    "Pioneered zero-shot CLIP use, but shows prompt engineering is brittle and needs improvement.", ACCENT_TEAL)

create_paper_slide(prs, "AnomalyCLIP: Object-Agnostic Prompt Learning for Zero-Shot AD", "Paper 3 | Zhou et al. [3] | ICLR 2024",
    ["Solves WinCLIP's manual prompt limitation", "Learns structured, object-agnostic prompt embeddings", "Captures universal 'normal' vs 'abnormal' semantics", "Frozen CLIP backbone; generalizes to unseen categories"],
    "~94.0% Image AUROC \nBetter cross-domain generalization.",
    ["Requires intensive prompt training phase", "Learned prompts aren't human-readable", "Still lacks natural language explanations"],
    "Automated prompts perform better, validating that prompt design is key to CLIP performance.", ACCENT_VIOLET)

create_paper_slide(prs, "AnomalyGPT: Industrial AD Using Large Vision-Language Models", "Paper 4 | Gu et al. [4] | AAAI 2024",
    ["Fine-tunes a multimodal LVLM (Vicuna) for defect detection", "Uses synthetic anomaly simulation (NSA) for training data", "Features an image decoder and visual-language prompt learner", "Supports multi-turn interactive dialogue for inspection"],
    "94.1% AUROC (1-shot) \nProvides conversational defect explanations.",
    ["Not zero-shot (needs fine-tuning)", "Relies on synthetic data quality", "Computationally expensive"],
    "First to provide LLM-based explanations, but loses the zero-shot magic by requiring fine-tuning.", ACCENT_INDIGO)

# ===== SLIDE 9: MMAD =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
create_header(s, "MMAD: Multi-Modal Anomaly Detection Benchmark", "Paper 5 | Jiang et al. [5] | ICLR 2025", ACCENT_AMBER)

add_rect(s, 0.8, 1.8, 5.5, 2.5, BG_CARD)
add_text(s, 1.1, 2.0, 5.1, 0.4, "Benchmark Scale", font_size=18, color=ACCENT_AMBER, bold=True)
add_bullets(s, 1.1, 2.5, 5.1, 1.5, [
    "39,672 queries across 8,366 images",
    "7 subtasks: classification, localization, type, severity, description, cause, repair",
    "Tested GPT-4o, Claude, LLaVA, Gemini"
], font_size=15)

add_rect(s, 7.0, 1.8, 5.5, 2.5, BG_CARD)
add_text(s, 7.3, 2.0, 5.1, 0.4, "Model Performance", font_size=18, color=ACCENT_ROSE, bold=True)
add_bullets(s, 7.3, 2.5, 5.1, 1.5, [
    "GPT-4o (Best): Only ~75% accuracy",
    "Open-source MLLMs perform poorly on fine-grained defects",
    "Models struggle with anomaly semantics vs standard objects"
], font_size=15)

add_rect(s, 0.8, 4.8, 11.7, 1.8, BG_CARD)
add_text(s, 1.1, 5.0, 11.1, 0.4, "Strategic Takeaway", font_size=18, color=ACCENT_TEAL, bold=True)
add_text(s, 1.1, 5.5, 11.1, 0.8, "Current Multimodal LLMs cannot perform reliable industrial inspection out-of-the-box. A hybrid approach utilizing specialized visual scoring (CLIP) paired with MLLM reasoning is required.", font_size=16, color=TEXT_MAIN, bold=True, align=PP_ALIGN.CENTER)

# ===== SLIDE 10: COMPARISON =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
create_header(s, "Comparative Analysis", "Synthesizing the Literature", ACCENT_INDIGO)

add_rect(s, 0.8, 1.8, 11.7, 0.6, ACCENT_INDIGO)
cols = [("Method", 1.0, 2.0), ("Venue", 3.2, 1.2), ("Zero-Shot?", 4.7, 1.3), ("AUROC", 6.2, 1.2), ("Explains?", 7.7, 1.2), ("Key Limitation", 9.2, 3.2)]
for text, left, w in cols:
    add_text(s, left, 1.9, w, 0.4, text, font_size=15, color=TEXT_MAIN, bold=True)

rows = [
    ("PatchCore", "CVPR'22", "No", "99.1%", "No", "Needs training per category"),
    ("WinCLIP", "CVPR'23", "Yes", "91.8%", "No", "Hand-crafted ad-hoc prompts"),
    ("AnomalyCLIP", "ICLR'24", "Yes*", "~94.0%", "No", "Needs prompt training phase"),
    ("AnomalyGPT", "AAAI'24", "No", "94.1%", "Yes", "Requires LVLM fine-tuning"),
]

for i, (m, v, zs, au, ex, lim) in enumerate(rows):
    y = 2.6 + i * 0.7
    add_rect(s, 0.8, y, 11.7, 0.6, BG_CARD if i % 2 == 0 else BG_BASE)
    c1 = ACCENT_TEAL if "CLIP" in m else TEXT_MAIN
    cx = ACCENT_ROSE if zs == "No" else ACCENT_TEAL
    ce = ACCENT_TEAL if ex == "Yes" else ACCENT_ROSE
    add_text(s, 1.0, y+0.1, 2.0, 0.4, m, font_size=15, color=c1, bold=True)
    add_text(s, 3.2, y+0.1, 1.2, 0.4, v, font_size=14, color=TEXT_MUTED)
    add_text(s, 4.7, y+0.1, 1.3, 0.4, zs, font_size=14, color=cx, bold=True)
    add_text(s, 6.2, y+0.1, 1.2, 0.4, au, font_size=15, color=TEXT_MAIN, bold=True)
    add_text(s, 7.7, y+0.1, 1.2, 0.4, ex, font_size=14, color=ce, bold=True)
    add_text(s, 9.2, y+0.1, 3.2, 0.4, lim, font_size=13, color=TEXT_MUTED)

add_rect(s, 0.8, 5.8, 11.7, 0.8, ACCENT_ROSE)
add_text(s, 1.0, 6.0, 11.3, 0.5, "Conclusion: No existing method achieves Zero-Shot + Explanations + Open-Source concurrently.", font_size=16, color=TEXT_MAIN, bold=True, align=PP_ALIGN.CENTER)

# ===== SLIDE 11: RESEARCH GAP =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
create_header(s, "Identified Research Gaps", None, ACCENT_ROSE)

gaps = [
    ("01", "Lack of Systematic Prompt Engineering", "WinCLIP uses ad-hoc text prompts. What grammatical structures, state words, and compositional ensembles actually work best for industrial anomaly detection?", ACCENT_INDIGO),
    ("02", "Untapped Open-Source MLLMs in Zero-Shot", "Proprietary models (GPT-4V) are too slow/costly. Fine-tuned models (AnomalyGPT) lose zero-shot advantages. Leveraging raw LLaVA is unexplored.", ACCENT_TEAL),
    ("03", "No Unified Hybrid Approach", "Models currently either score well (CLIP) or reason well (MLLM). Combining CLIP's robust localization precision with MLLM's reasoning engine in a unified, training-free pipeline.", ACCENT_VIOLET)
]

for i, (num, title, desc, clr) in enumerate(gaps):
    y = 1.8 + i * 1.6
    add_rect(s, 0.8, y, 11.7, 1.3, BG_CARD)
    add_rect(s, 0.8, y, 0.1, 1.3, clr, shape_type=MSO_SHAPE.RECTANGLE)
    add_text(s, 1.2, y+0.15, 1.0, 0.6, num, font_size=32, color=clr, bold=True, font=FONT_TITLE)
    add_text(s, 2.2, y+0.15, 10.0, 0.4, title, font_size=18, color=TEXT_MAIN, bold=True)
    add_text(s, 2.2, y+0.6, 10.0, 0.6, desc, font_size=14, color=TEXT_MUTED)

# ===== SLIDE 12: PROBLEM STATEMENT (MATHEMATICAL) =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
create_header(s, "Formal Problem Statement", "Mathematical Formulation of the Proposed Architecture", ACCENT_INDIGO)

# Overall Definition
add_rect(s, 0.8, 1.8, 11.7, 1.0, BG_CARD)
add_rect(s, 0.8, 1.8, 0.1, 1.0, ACCENT_TEAL, shape_type=MSO_SHAPE.RECTANGLE)
add_text(s, 1.1, 2.0, 11.0, 0.8, "Design and development of a zero-shot anomaly detection and reasoning framework using Vision-Language Models (VLM) for precise localization and Multimodal Large Language Models (MLLM) for comprehensive semantic explanation.", font_size=16, color=TEXT_MAIN, bold=True)

# Math Formulation (VLM Stage)
add_rect(s, 0.8, 3.2, 5.5, 2.5, BG_CARD)
add_text(s, 1.1, 3.4, 5.0, 0.3, "Stage 1: VLM Scoring  S_vlm(x)", font_size=18, color=ACCENT_INDIGO, bold=True)
add_bullets(s, 1.1, 3.9, 5.0, 1.8, [
    "Given test image x ∈ ℝ^(H×W×3)",
    "Normality score f(x) computed via cosine similarity:",
    "  S_vlm(x) = 1 - max_i [ (v · t_i) / (||v|| ||t_i||) ]",
    "Where v is image feature, t_i are text prompts",
    "Outputs: Anomaly score S_vlm(x) and heatmap M(x)"
], font_size=14)

# Math Formulation (MLLM Stage)
add_rect(s, 7.0, 3.2, 5.5, 2.5, BG_CARD)
add_text(s, 7.3, 3.4, 5.0, 0.3, "Stage 2: MLLM Reasoning  R_mllm(x, M)", font_size=18, color=ACCENT_VIOLET, bold=True)
add_bullets(s, 7.3, 3.9, 5.0, 1.8, [
    "Conditional Execution:",
    "  Let decision D(x) = 1 if S_vlm(x) > τ, else 0",
    "If D(x) == 1:",
    "  R_mllm = MLLM(x, M(x), P_query)",
    "Where P_query asks for semantic reasoning of M(x)"
], font_size=14)

add_rect(s, 0.8, 6.0, 11.7, 0.8, ACCENT_TEAL)
add_text(s, 1.0, 6.2, 11.3, 0.5, "Objective: Maximize detection accuracy P(D(x)=y) while generating human-readable explanation R.", font_size=16, color=TEXT_MAIN, bold=True, align=PP_ALIGN.CENTER)

# ===== SLIDE 13: METHODOLOGY =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
create_header(s, "Proposed Research Methodology", "Two-Stage Zero-Shot Interpretive Architecture", ACCENT_TEAL)

# Stage 1
add_rect(s, 0.8, 2.0, 5.0, 3.8, BG_CARD)
add_rect(s, 0.8, 2.0, 5.0, 0.6, ACCENT_INDIGO, shape_type=MSO_SHAPE.RECTANGLE)
add_text(s, 1.0, 2.1, 4.6, 0.4, "Stage 1: VLM Scoring (CLIP)", font_size=18, color=TEXT_MAIN, bold=True)
add_bullets(s, 1.1, 2.8, 4.5, 2.5, [
    "Pre-trained CLIP (Frozen)",
    "Systematic Prompt Array vs Visual Patches",
    "Generates Anomaly Map & AUROC Score",
    "Acts as an ultra-fast filter"
], font_size=15)

# Arrow
add_text(s, 6.0, 3.6, 1.0, 0.5, "➔", font_size=40, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

# Stage 2
add_rect(s, 7.3, 2.0, 5.2, 3.8, BG_CARD)
add_rect(s, 7.3, 2.0, 5.2, 0.6, ACCENT_VIOLET, shape_type=MSO_SHAPE.RECTANGLE)
add_text(s, 7.5, 2.1, 4.8, 0.4, "Stage 2: MLLM Reasoning (LLaVA)", font_size=18, color=TEXT_MAIN, bold=True)
add_bullets(s, 7.5, 2.8, 4.8, 2.5, [
    "Triggered ONLY if Stage 1 Score > Threshold",
    "Analyzes defect using original image + heatmap",
    "Outputs localized bounding boxes",
    "Generates natural language repair report"
], font_size=15)

add_rect(s, 2.0, 6.2, 9.333, 0.8, ACCENT_TEAL)
add_text(s, 2.0, 6.4, 9.333, 0.4, "Zero-Shot Precision  +  Interpretability  +  Open-Source Scalability", font_size=18, color=BG_BASE, bold=True, align=PP_ALIGN.CENTER)

# ===== SLIDE 13: CONTRIBUTIONS =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
create_header(s, "Expected Project Contributions", None, ACCENT_VIOLET)

for i, (label, text, details) in enumerate([
    ("1", "Comprehensive Prompt Engineering Framework", "Empirical study on grammatical structure and ensemble design for CLIP in manufacturing contexts."),
    ("2", "Novel Hybrid VLM+MLLM Pipeline", "A cohesive, non-fine-tuned architecture blending CLIP localization with LLaVA linguistic reasoning."),
    ("3", "Robust Empirical Benchmarking", "Evaluation on MVTec AD and VisA vs SOTA baselines (PatchCore, WinCLIP) to prove efficacy and limits.")
]):
    y = 1.8 + i * 1.6
    add_rect(s, 0.8, y, 11.7, 1.3, BG_CARD)
    # Circle indicator
    add_rect(s, 1.2, y+0.25, 0.8, 0.8, ACCENT_VIOLET, shape_type=MSO_SHAPE.OVAL)
    add_text(s, 1.2, y+0.4, 0.8, 0.8, label, font_size=24, color=TEXT_MAIN, bold=True, align=PP_ALIGN.CENTER)
    
    add_text(s, 2.5, y+0.2, 9.5, 0.4, text, font_size=18, color=ACCENT_TEAL, bold=True)
    add_text(s, 2.5, y+0.65, 9.5, 0.5, details, font_size=15, color=TEXT_MUTED)

# ===== SLIDE 14: TIMELINE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
create_header(s, "Project Timeline & Execution Plan", "Strategic Delivery Roadmap", ACCENT_INDIGO)

timings = [
    ("Phase 1: Foundation", "Jan - May 2026", "Lit survey, environment prep, and initial CLIP staging", ACCENT_INDIGO),
    ("Phase 2: Core Development", "Jun - Oct 2026", "Prompt engineering study & Hybrid pipeline assembly", ACCENT_TEAL),
    ("Phase 3: Validation", "Nov - Jan 2027", "Ablation profiling on MVTec AD/VisA environments", ACCENT_VIOLET),
    ("Phase 4: Finalization", "Feb - May 2027", "Thesis authoring and potential conference submission", ACCENT_AMBER)
]

for i, (p, period, task, clr) in enumerate(timings):
    x = 0.8 + (i * 3.0)
    add_rect(s, x, 2.0, 2.7, 3.8, BG_CARD)
    add_rect(s, x, 2.0, 2.7, 0.2, clr, shape_type=MSO_SHAPE.RECTANGLE)
    add_text(s, x+0.1, 2.4, 2.5, 0.4, period, font_size=14, color=TEXT_MUTED, bold=True)
    add_text(s, x+0.1, 2.8, 2.5, 0.5, p, font_size=17, color=clr, bold=True)
    add_text(s, x+0.1, 3.5, 2.5, 2.0, task, font_size=15, color=TEXT_MAIN)

add_rect(s, 0.8, 6.2, 11.7, 0.8, BG_CARD)
add_text(s, 1.0, 6.45, 11.3, 0.4, "Tech Stack: PyTorch, HuggingFace, Anomalib, OpenCV, MVTec AD", font_size=15, color=TEXT_MUTED, align=PP_ALIGN.CENTER, bold=True)

# ===== SLIDE 15: THANK YOU =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
add_rect(s, 4, 2.5, 5.333, 2.5, BG_CARD)
add_line(s, 4, 2.5, 5.333, ACCENT_TEAL)
add_text(s, 4, 3.0, 5.333, 1.0, "Thank You", font_size=42, color=TEXT_MAIN, bold=True, align=PP_ALIGN.CENTER, font=FONT_TITLE)
add_text(s, 4, 4.0, 5.333, 0.5, "Questions & Feedback Welcome", font_size=18, color=ACCENT_TEAL, align=PP_ALIGN.CENTER)

add_text(s, 1, 6.0, 11.333, 0.5, "Reach Out: [email@nitc.ac.in]  •  [Your Name]", font_size=14, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

# ===== SLIDE 17: REFERENCES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(s)
create_header(s, "References", "Literature Survey Sources", ACCENT_INDIGO)

refs = [
    "[1] K. Roth et al., 'Towards Total Recall in Industrial Anomaly Detection,' CVPR, 2022.",
    "[2] J. Jeong et al., 'WinCLIP: Zero-/Few-Shot Anomaly Classification and Segmentation,' CVPR, 2023.",
    "[3] Q. Zhou et al., 'AnomalyCLIP: Object-Agnostic Prompt Learning for Zero-Shot Anomaly Detection,' ICLR, 2024.",
    "[4] Z. Gu et al., 'AnomalyGPT: Detecting Industrial Anomalies Using Large Vision-Language Models,' AAAI, 2024.",
    "[5] X. Jiang et al., 'MMAD: A Comprehensive Benchmark for Multimodal Large Language Models in Industrial Anomaly Detection,' ICLR, 2025."
]

add_bullets(s, 1.0, 2.0, 11.0, 4.0, refs, font_size=14, color=TEXT_MAIN)

# ===== SAVE =====
out_path = r"p:\Research\multimodal-anomaly-detection\docs\proposal\Literature_Survey_Presentation_v2.pptx"
prs.save(out_path)
print(f"Presentation generated successfully: {out_path}")
print(f"Total Slides: {len(prs.slides)}")
